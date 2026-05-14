"""
gen_pptx.py — SO-ARM101 LeRobot 5-page Beamer PPTX replicator
Tsinghua purple theme, light background.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

BASE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(BASE, "sources")

# ── Tsinghua Colors ──
THUPURPLE = RGBColor(102, 45, 145)
THUGOLD   = RGBColor(166, 130, 46)
LGRAY     = RGBColor(245, 245, 248)
WHITE     = RGBColor(255, 255, 255)
GRAY      = RGBColor(120, 120, 120)
DARKTEXT  = RGBColor(40, 40, 40)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_rounded_rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def add_textbox(slide, left, top, w, h, text, size=12, bold=False,
                color=DARKTEXT, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(size); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font; p.alignment = align
    return tb

def add_title_bar(slide, title):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.55), THUPURPLE)
    add_textbox(slide, Inches(0.5), Inches(0.07), Inches(12), Inches(0.45),
                title, size=20, bold=True, color=WHITE)
    add_rect(slide, Inches(0), Inches(0.55), SLIDE_W, Inches(0.04), THUGOLD)

def add_block(slide, left, top, w, h, title, items, title_color=THUPURPLE):
    bh = Inches(0.30)
    add_rounded_rect(slide, left, top, w, bh, title_color)
    add_textbox(slide, left+Inches(0.12), top+Inches(0.02), w-Inches(0.24), bh,
                title, size=10, bold=True, color=WHITE)
    bt = top + bh; body_h = h - bh
    add_rounded_rect(slide, left, bt, w, body_h, LGRAY)
    y = bt + Inches(0.06)
    for item in items:
        add_textbox(slide, left+Inches(0.18), y, w-Inches(0.36), Inches(0.22),
                    item, size=9, color=DARKTEXT)
        y += Inches(0.22)

def add_example_block(slide, left, top, w, h, title, text):
    bh = Inches(0.28)
    add_rounded_rect(slide, left, top, w, bh, RGBColor(140, 108, 28))
    add_textbox(slide, left+Inches(0.1), top+Inches(0.02), w-Inches(0.2), bh,
                title, size=9, bold=True, color=WHITE)
    bt = top + bh
    add_rounded_rect(slide, left, bt, w, h - bh, RGBColor(250, 248, 240))
    add_textbox(slide, left+Inches(0.1), bt+Inches(0.04), w-Inches(0.2), h-bh-Inches(0.08),
                text, size=9, color=DARKTEXT)

def add_alert_block(slide, left, top, w, h, title, text):
    bh = Inches(0.28)
    add_rounded_rect(slide, left, top, w, bh, RGBColor(180, 50, 50))
    add_textbox(slide, left+Inches(0.1), top+Inches(0.02), w-Inches(0.2), bh,
                title, size=9, bold=True, color=WHITE)
    bt = top + bh
    add_rounded_rect(slide, left, bt, w, h - bh, RGBColor(255, 240, 240))
    add_textbox(slide, left+Inches(0.1), bt+Inches(0.04), w-Inches(0.2), h-bh-Inches(0.08),
                text, size=9, color=DARKTEXT)

def add_image_safe(slide, filename, left, top, width=None, height=None):
    path = os.path.join(SRC, filename)
    if not os.path.exists(path):
        add_textbox(slide, left, top, Inches(3), Inches(0.4),
                    f"[Image not found: {filename}]", size=10, color=RGBColor(200,0,0))
        return
    kw = {"image_file": path, "left": left, "top": top}
    if width: kw["width"] = width
    if height: kw["height"] = height
    slide.shapes.add_picture(**kw)

def add_page_number(slide, n, total):
    add_textbox(slide, Inches(12.2), Inches(7.05), Inches(1), Inches(0.3),
                f"{n}/{total}", size=9, color=GRAY, align=PP_ALIGN.RIGHT)

def add_table(slide, left, top, col_widths, headers, rows, hdr_color=THUPURPLE):
    rh = Inches(0.28)
    x = left
    for hdr, cw in zip(headers, col_widths):
        add_rounded_rect(slide, x, top, cw, rh, hdr_color)
        add_textbox(slide, x+Inches(0.04), top+Inches(0.02), cw-Inches(0.08), rh,
                    hdr, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw
    top += rh
    for ri, row in enumerate(rows):
        x = left
        bg = LGRAY if ri % 2 == 0 else WHITE
        for val, cw in zip(row, col_widths):
            add_rounded_rect(slide, x, top, cw, rh, bg)
            add_textbox(slide, x+Inches(0.04), top+Inches(0.02), cw-Inches(0.08), rh,
                        str(val), size=8, color=DARKTEXT, align=PP_ALIGN.CENTER)
            x += cw
        top += rh
    return top

# ════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
TOTAL = 5

# ── Slide 1: Title (white bg) ──
s = add_blank_slide(prs)
add_textbox(s, Inches(1), Inches(2.0), Inches(11.333), Inches(0.9),
            "SO-101 Training Framework", size=32, bold=True, color=THUPURPLE,
            align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(2.9), Inches(11.333), Inches(0.5),
            "& Orchestrator Architecture", size=18, color=RGBColor(140, 90, 180),
            align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(4.0), Inches(11.333), Inches(0.35),
            "Phang Hong Hao (潘洪浩)", size=16, color=DARKTEXT, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(4.5), Inches(11.333), Inches(0.3),
            "Tsinghua University", size=12, color=GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, Inches(1), Inches(4.9), Inches(11.333), Inches(0.3),
            "2026-05-13", size=11, color=RGBColor(180,180,180), align=PP_ALIGN.CENTER)

# ── Slide 2: Training Framework & Speed ──
s = add_blank_slide(prs)
add_title_bar(s, "Training Framework & Speed Benchmark")

add_textbox(s, Inches(0.5), Inches(0.75), Inches(6), Inches(0.3),
            "Two Policies on RTX 6000D", size=13, bold=True, color=THUPURPLE)

cw1 = [Inches(2.5), Inches(2.5)]
y = Inches(1.1)
y = add_table(s, Inches(0.5), y, cw1,
              ["PushT Diffusion", ""],
              [["Policy", "Diffusion (UNet)"],
               ["Params", "263M"],
               ["Batch", "64"],
               ["Steps", "100K"]])
y += Inches(0.15)
add_table(s, Inches(0.5), y, cw1,
          ["SO-101 ACT", ""],
          [["Policy", "ACT (Transformer)"],
           ["Params", "84M"],
           ["Batch", "128 (max)"],
           ["Steps", "50K"]])

add_textbox(s, Inches(7.0), Inches(0.75), Inches(6), Inches(0.3),
            "torchcodec vs pyav", size=13, bold=True, color=THUPURPLE)

cw2 = [Inches(3.0), Inches(1.6), Inches(1.2)]
add_table(s, Inches(7.0), Inches(1.1), cw2,
          ["Config", "Speed", "Factor"],
          [["PushT pyav", "0.5-0.7 step/s", "1x"],
           ["PushT torchcodec", "12-15 step/s", "20x"],
           ["ACT pyav 8w", "0.38 step/s", "1x"],
           ["ACT pyav 8w AMP", "0.45 step/s", "1.2x"],
           ["ACT torchcodec 8w AMP", "1.9 step/s", "8.6x"]])

add_example_block(s, Inches(7.0), Inches(4.2), Inches(5.7), Inches(0.85),
                  "Key Finding",
                  "GPU 100%, 视频解码是真正瓶颈。torchcodec GPU 硬件加速 >> pyav CPU 逐帧解码。")

add_page_number(s, 2, TOTAL)

# ── Slide 3: Optimal Config & Bugs ──
s = add_blank_slide(prs)
add_title_bar(s, "Optimal Configuration & Known Bugs")

add_textbox(s, Inches(0.5), Inches(0.75), Inches(6), Inches(0.3),
            "Launch Command", size=13, bold=True, color=THUPURPLE)

cw3 = [Inches(2.2), Inches(2.5)]
add_table(s, Inches(0.5), Inches(1.1), cw3,
          ["Parameter", "Value"],
          [["video_backend", "torchcodec"],
           ["policy.type", "act"],
           ["policy.device", "cuda (not cuda:N!)"],
           ["chunk_size", "100"],
           ["dim_model", "512"],
           ["use_amp", "True"],
           ["batch_size", "128 (max)"],
           ["num_workers", "8"]])

add_alert_block(s, Inches(0.5), Inches(4.5), Inches(5.0), Inches(0.8),
                "Why batch_size=128 Max?",
                "ACT chunk_size=100 O(n²), 128 batch → ~79GB VRAM (RTX 6000D limit)")

add_textbox(s, Inches(7.0), Inches(0.75), Inches(5.5), Inches(0.3),
            "Known Bugs & Workarounds", size=13, bold=True, color=THUPURPLE)

cw4 = [Inches(2.5), Inches(3.2)]
add_table(s, Inches(7.0), Inches(1.1), cw4,
          ["Bug", "Fix"],
          [["is_amp_available ValueError", "CUDA_VISIBLE_DEVICES=6"],
           ["torchcodec DecodingError", "同上"],
           ["CXXABI_1.3.15 not found", "LD_PRELOAD=.../libstdc++.so.6"],
           ["libavutil.so.X not found", "conda install ffmpeg"]])

add_page_number(s, 3, TOTAL)

# ── Slide 4: Orchestrator Architecture ──
s = add_blank_slide(prs)
add_title_bar(s, "Orchestrator Architecture")

add_image_safe(s, "orchestrator_pipeline.png",
               Inches(1.5), Inches(0.75), width=Inches(10.3))

y = Inches(4.5)
add_block(s, Inches(0.5), y, Inches(5.8), Inches(2.2),
          "Core Modules",
          ["• arm101_orchestrator.py — 主控 + CLI",
           "• phase_manager.py — YAML 三层合并",
           "• state_store.py — 崩溃恢复 (原子 JSON)",
           "• data_collector.py — MuJoCo 采集"])

add_block(s, Inches(6.8), y, Inches(5.8), Inches(2.2),
          "Key Features",
          ["• --fresh / --start-from <id> / --dry-run",
           "• 崩溃恢复: 重连训练 PID",
           "• Loss Monitor: CONVERGED / OVERFITTING"])

add_page_number(s, 4, TOTAL)

# ── Slide 5: Reference Models & Takeaways ──
s = add_blank_slide(prs)
add_title_bar(s, "Reference Models & Key Takeaways")

add_textbox(s, Inches(0.5), Inches(0.75), Inches(6), Inches(0.3),
            "HuggingFace Models (Survey 100+)", size=13, bold=True, color=THUPURPLE)

cw5 = [Inches(0.5), Inches(5.0)]
add_table(s, Inches(0.5), Inches(1.1), cw5,
          ["#", "Model / Key Info"],
          [["1", "Sa74ll/smolvla_so101 — SmolVLA, 87.66% 成功率"],
           ["2", "TakuyaHiraoka/act_so101 — ACT, 真实机械臂"],
           ["3", "davidlinjiahao/lerobot_so101 — ACT, MuJoCo, 最相似"]])

add_example_block(s, Inches(0.5), Inches(3.1), Inches(5.8), Inches(0.8),
                  "Pipeline",
                  "scene.xml → MjModel → trajectory → render → LeRobotDataset → HF Hub")

add_textbox(s, Inches(7.0), Inches(0.75), Inches(5.5), Inches(0.3),
            "Key Takeaways", size=13, bold=True, color=THUPURPLE)

items = [
    ("1. ACT 是 SO-101 主流策略", "HF 上 30+ ACT 模型"),
    ("2. SmolVLA 成功率更高", "Sa74ll 达 87.66%，靠数据分割"),
    ("3. 本项目与 davidlinjiahao 最相似", "都是 MuJoCo + ACT + SO-101"),
    ("4. 视频解码是训练瓶颈", "torchcodec 比 pyav 快 8-20x"),
]
y = Inches(1.15)
for title, desc in items:
    add_textbox(s, Inches(7.2), y, Inches(5.3), Inches(0.25),
                title, size=11, bold=True, color=DARKTEXT)
    add_textbox(s, Inches(7.5), y + Inches(0.25), Inches(5.0), Inches(0.22),
                f"• {desc}", size=9, color=GRAY)
    y += Inches(0.55)

add_page_number(s, 5, TOTAL)

# ── Save ──
out = os.path.join(BASE, "main.pptx")
try:
    prs.save(out)
    print(f"Saved: {out}")
except PermissionError:
    out2 = os.path.join(BASE, "main_v2.pptx")
    prs.save(out2)
    print(f"File locked, saved as: {out2}")
